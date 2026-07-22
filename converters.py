import os
import io
import tempfile
import subprocess
from PIL import Image
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches

def resize_and_pad(image, target_size=(1920, 1080), background_color=(255, 255, 255)):
    """
    Fits a PIL image into a target canvas size, letterboxing/centering it.
    """
    # Create a copy so we don't modify the original
    img = image.copy()
    img.thumbnail(target_size, Image.Resampling.LANCZOS)
    
    # Create background canvas
    canvas = Image.new("RGB", target_size, background_color)
    
    # Paste centered
    x = (target_size[0] - img.width) // 2
    y = (target_size[1] - img.height) // 2
    
    # If image has transparency (RGBA/LA) or is palette with transparency
    if img.mode in ('RGBA', 'LA') or (img.mode == 'P' and 'transparency' in img.info):
        canvas.paste(img, (x, y), img.convert("RGBA"))
    else:
        canvas.paste(img, (x, y))
        
    return canvas

def pdf_to_png(pdf_path, output_dir, res_mode="height_1080", scale_mode="fit"):
    """
    Renders PDF pages to PNG images at selected resolution.
    Returns list of paths to generated images.
    """
    doc = fitz.open(pdf_path)
    generated_files = []
    
    os.makedirs(output_dir, exist_ok=True)
    
    for i, page in enumerate(doc):
        rect = page.rect
        width, height = rect.width, rect.height
        
        # Calculate matrix scaling factor based on resolution mode
        if res_mode == "height_1080":
            zoom = 1080 / height
        elif res_mode == "width_1920":
            zoom = 1920 / width
        elif res_mode == "fit_16_9":
            # Scale to fit inside 1920x1080
            zoom = min(1920 / width, 1080 / height)
        elif res_mode == "original":
            zoom = 2.0  # Render at 2x default resolution (144 DPI) for clarity
        else:
            zoom = 1080 / height
            
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        
        out_name = f"page_{i+1:03d}.png"
        out_path = os.path.join(output_dir, out_name)
        
        if res_mode == "fit_16_9" and scale_mode == "fit":
            # Render page first, then use PIL to letterbox it
            img_data = pix.tobytes("png")
            img = Image.open(io.BytesIO(img_data))
            padded_img = resize_and_pad(img, (1920, 1080))
            padded_img.save(out_path, "PNG")
        else:
            pix.save(out_path)
            
        generated_files.append(out_path)
        
    doc.close()
    return generated_files

def pdf_to_pptx(pdf_path, output_pptx_path, res_mode="fit_16_9", temp_dir=None):
    """
    Converts PDF to PPTX by first converting PDF pages to PNGs, then adding them to PPTX.
    Satisfies user request: "from pdf to ppt, but converts first to png and then creates ppt out of png:s"
    """
    if temp_dir is None:
        # Create a nested temp directory specifically for this conversion
        temp_dir_obj = tempfile.TemporaryDirectory()
        temp_dir = temp_dir_obj.name
    else:
        temp_dir_obj = None

    try:
        # 1. Convert pages to PNG first (with 1080 height/fit for high fidelity)
        png_paths = pdf_to_png(pdf_path, temp_dir, res_mode=res_mode, scale_mode="fit")
        
        # 2. Build PPTX from these images
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 Aspect Ratio
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        
        for png_path in png_paths:
            slide = prs.slides.add_slide(blank_layout)
            
            img = Image.open(png_path)
            img_w, img_h = img.size
            img_aspect = img_w / img_h
            
            slide_aspect = 13.333 / 7.5
            
            if abs(img_aspect - slide_aspect) < 0.02:
                # Fill screen if aspect matches 16:9
                slide.shapes.add_picture(png_path, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
            else:
                # Center and fit (letterbox) on the slide
                if img_aspect > slide_aspect:
                    # Width limited
                    w = prs.slide_width
                    h = prs.slide_width / img_aspect
                    left = Inches(0)
                    top = (prs.slide_height - h) / 2
                else:
                    # Height limited
                    h = prs.slide_height
                    w = prs.slide_height * img_aspect
                    left = (prs.slide_width - w) / 2
                    top = Inches(0)
                slide.shapes.add_picture(png_path, left, top, width=w, height=h)
            img.close()
            
        prs.save(output_pptx_path)
    finally:
        if temp_dir_obj:
            temp_dir_obj.cleanup()

def png_to_pdf(img_paths, output_pdf_path, res_mode="height_1080"):
    """
    Converts a list of images into a single PDF document.
    """
    doc = fitz.open()
    for img_path in img_paths:
        img = Image.open(img_path)
        w, h = img.size
        
        if res_mode == "height_1080":
            target_h = 1080
            target_w = int(w * (1080 / h))
            resized_img = img.resize((target_w, target_h), Image.Resampling.LANCZOS)
        elif res_mode == "width_1920":
            target_w = 1920
            target_h = int(h * (1920 / w))
            resized_img = img.resize((target_w, target_h), Image.Resampling.LANCZOS)
        elif res_mode == "fit_16_9":
            target_w = 1920
            target_h = 1080
            resized_img = resize_and_pad(img, (1920, 1080))
        elif res_mode == "original":
            target_w = w
            target_h = h
            resized_img = img.copy()
        else:
            target_h = 1080
            target_w = int(w * (1080 / h))
            resized_img = img.resize((target_w, target_h), Image.Resampling.LANCZOS)
            
        img_byte_arr = io.BytesIO()
        resized_img.save(img_byte_arr, format='PNG')
        img_bytes = img_byte_arr.getvalue()
        
        # Create PDF page from raw PNG bytes
        img_pdf_bytes = fitz.open("png", img_bytes).convert_to_pdf()
        img_pdf = fitz.open("pdf", img_pdf_bytes)
        
        page = doc.new_page(width=target_w, height=target_h)
        page.show_pdf_page(page.rect, img_pdf, 0)
        
        img.close()
        resized_img.close()
        img_pdf.close()
        
    doc.save(output_pdf_path)
    doc.close()

def png_to_pptx(img_paths, output_pptx_path, res_mode="fit_16_9"):
    """
    Converts list of images to a single 1080p PPTX presentation.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    for img_path in img_paths:
        slide = prs.slides.add_slide(blank_layout)
        
        img = Image.open(img_path)
        img_w, img_h = img.size
        img_aspect = img_w / img_h
        
        slide_aspect = 13.333 / 7.5
        
        if res_mode == "fit_16_9":
            if img_aspect > slide_aspect:
                # Width limited
                w = prs.slide_width
                h = prs.slide_width / img_aspect
                left = Inches(0)
                top = (prs.slide_height - h) / 2
            else:
                # Height limited
                h = prs.slide_height
                w = prs.slide_height * img_aspect
                left = (prs.slide_width - w) / 2
                top = Inches(0)
        else:
            # Stretch/Fill slide
            w = prs.slide_width
            h = prs.slide_height
            left = Inches(0)
            top = Inches(0)
            
        slide.shapes.add_picture(img_path, left, top, width=w, height=h)
        img.close()
        
    prs.save(output_pptx_path)

def pptx_to_pdf(pptx_path, output_pdf_path):
    """
    Converts a PPTX presentation to PDF using Microsoft PowerPoint via AppleScript and macOS launch services.
    Bypasses macOS App Sandbox file permission errors.
    """
    import time
    abs_pptx = os.path.abspath(pptx_path)
    abs_pdf = os.path.abspath(output_pdf_path)
    
    if os.path.exists(abs_pdf):
        os.remove(abs_pdf)
        
    try:
        # 1. Open PowerPoint file in background via launch services (resolves sandbox access automatically)
        subprocess.run(['open', '-g', '-a', 'Microsoft PowerPoint', abs_pptx])
        
        # Wait up to 5 seconds for active presentation to load, polling every 0.5s
        loaded = False
        for _ in range(10):
            # Check if there is an active presentation
            check_script = 'tell application "Microsoft PowerPoint" to get name of active presentation'
            res = subprocess.run(['osascript', '-e', check_script], stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
            if res.returncode == 0:
                loaded = True
                break
            time.sleep(0.5)
            
        if not loaded:
            raise RuntimeError("Microsoft PowerPoint failed to open the presentation in a timely manner.")
            
        # 2. Save as PDF via AppleScript
        save_script = f'tell application "Microsoft PowerPoint" to save active presentation in (POSIX file "{abs_pdf}") as save as PDF'
        proc_save = subprocess.run(
            ['osascript', '-e', save_script],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True
        )
        
        if proc_save.returncode != 0:
            raise RuntimeError(f"PowerPoint save-to-PDF failed: {proc_save.stderr or proc_save.stdout}")
            
    finally:
        # 3. Close the presentation without saving changes
        close_script = 'tell application "Microsoft PowerPoint" to close active presentation saving no'
        subprocess.run(['osascript', '-e', close_script], stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
    if not os.path.exists(abs_pdf):
        raise RuntimeError("PowerPoint-to-PDF conversion finished, but the output PDF was not found.")
        
    return abs_pdf

def pptx_to_png(pptx_path, output_dir, res_mode="height_1080", temp_dir=None):
    """
    Converts PPTX slides to PNG images by converting to PDF first, then rendering PDF to PNG.
    """
    if temp_dir is None:
        temp_dir_obj = tempfile.TemporaryDirectory()
        temp_dir = temp_dir_obj.name
    else:
        temp_dir_obj = None
        
    try:
        temp_pdf_path = os.path.join(temp_dir, "temp_presentation.pdf")
        
        # 1. PPTX -> PDF via AppleScript
        pptx_to_pdf(pptx_path, temp_pdf_path)
        
        # 2. PDF -> PNG via PyMuPDF
        png_paths = pdf_to_png(temp_pdf_path, output_dir, res_mode=res_mode, scale_mode="fit")
        return png_paths
    finally:
        if temp_dir_obj:
            temp_dir_obj.cleanup()
